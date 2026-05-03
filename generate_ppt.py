"""
Generate MTech BITS Pilani Mid-Semester Dissertation PPT
for Corporate Car Pooling System
"""

import os
import io
import cairosvg
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy

# ─── Colour palette ────────────────────────────────────────────────────────────
PRIMARY    = RGBColor(0x1A, 0x37, 0x6C)   # deep navy
ACCENT     = RGBColor(0x0D, 0x9E, 0xDF)   # sky blue
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF4, 0xF8)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x44)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)

DIAGRAMS = "/home/user/Corporate-Car-Pooling-System/Diagrams"

# ─── Helpers ───────────────────────────────────────────────────────────────────

def svg_to_png_bytes(svg_path: str, scale: float = 2.0) -> bytes:
    with open(svg_path, "rb") as f:
        return cairosvg.svg2png(file_obj=f, scale=scale)


def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def fill_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_header_bar(slide, title_text, subtitle_text=None):
    """Dark navy top bar with title."""
    W = Inches(13.33)
    add_rect(slide, 0, 0, W, Inches(1.25), PRIMARY)
    add_text_box(slide, title_text,
                 Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.7),
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_text_box(slide, subtitle_text,
                     Inches(0.35), Inches(0.78), Inches(12.5), Inches(0.42),
                     font_size=14, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)


def add_footer(slide, text="MTech BITS Pilani | Corporate Car Pooling System | 2025"):
    W = Inches(13.33)
    add_rect(slide, 0, Inches(7.35), W, Inches(0.15), PRIMARY)
    add_text_box(slide, text,
                 Inches(0.25), Inches(7.37), Inches(12.8), Inches(0.25),
                 font_size=8, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)


def add_diagram_slide(prs, title, subtitle, svg_filename,
                      notes=None, img_top=Inches(1.35), img_height=Inches(5.8)):
    slide = add_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    add_header_bar(slide, title, subtitle)
    add_footer(slide)

    svg_path = os.path.join(DIAGRAMS, svg_filename)
    png_bytes = svg_to_png_bytes(svg_path)
    img_stream = io.BytesIO(png_bytes)

    # Calculate width to preserve aspect ratio
    from PIL import Image
    img = Image.open(io.BytesIO(png_bytes))
    orig_w, orig_h = img.size
    aspect = orig_w / orig_h
    img_w = img_height * aspect
    max_w = Inches(12.8)
    if img_w > max_w:
        img_w = max_w
        img_height = img_w / aspect

    left = (Inches(13.33) - img_w) / 2
    slide.shapes.add_picture(io.BytesIO(png_bytes), left, img_top, img_w, img_height)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def bullet_items(tf, items, font_size=16, color=DARK_GRAY, indent=0):
    """Add bullet items to an existing text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = indent
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=15, color=DARK_GRAY, bg=WHITE, border=None):
    shape = add_rect(slide, left, top, width, height, bg,
                     line_color=border or ACCENT, line_width=Pt(1.2))
    tf = shape.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return shape


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 – Title"""
    slide = add_slide(prs)
    fill_bg(slide, PRIMARY)

    # Large gradient-ish accent strip
    add_rect(slide, 0, Inches(4.6), Inches(13.33), Inches(0.08), ACCENT)

    # Logo/icon placeholder text
    add_text_box(slide, "🚗", Inches(5.9), Inches(0.55), Inches(1.5), Inches(1.2),
                 font_size=60, color=ACCENT, align=PP_ALIGN.CENTER)

    # Main title
    add_text_box(slide,
                 "Corporate Car Pooling System",
                 Inches(1.0), Inches(1.6), Inches(11.33), Inches(1.0),
                 font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide,
                 "Mid-Semester Dissertation Presentation",
                 Inches(1.0), Inches(2.55), Inches(11.33), Inches(0.6),
                 font_size=22, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "MTech (Software Engineering)  |  BITS Pilani",
                 Inches(1.0), Inches(3.15), Inches(11.33), Inches(0.5),
                 font_size=18, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

    # Presenter
    add_rect(slide, Inches(4.0), Inches(4.8), Inches(5.33), Inches(1.6),
             RGBColor(0x0D, 0x22, 0x4A))
    add_text_box(slide, "Presented by:",
                 Inches(4.15), Inches(4.9), Inches(5.0), Inches(0.35),
                 font_size=12, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Prashant Sopan Berad",
                 Inches(4.15), Inches(5.22), Inches(5.0), Inches(0.45),
                 font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "May 2025",
                 Inches(4.15), Inches(5.65), Inches(5.0), Inches(0.35),
                 font_size=13, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

    add_footer(slide)


def slide_agenda(prs):
    """Slide 2 – Agenda"""
    slide = add_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    add_header_bar(slide, "Agenda", "What we'll cover today")
    add_footer(slide)

    items = [
        ("01", "Problem Statement & Motivation"),
        ("02", "Proposed Solution & Objectives"),
        ("03", "System Architecture & Tech Stack"),
        ("04", "Core Modules Overview"),
        ("05", "Database Design (ER Diagram)"),
        ("06", "System Diagrams (Use Case, Activity, Sequence, State Machine)"),
        ("07", "Progress Made So Far"),
        ("08", "Backend & Frontend Implementation"),
        ("09", "AI Agent Module"),
        ("10", "Future Roadmap"),
    ]

    col_w = Inches(5.9)
    col_gap = Inches(0.5)
    left1 = Inches(0.35)
    left2 = left1 + col_w + col_gap
    top_start = Inches(1.4)
    row_h = Inches(0.57)

    for i, (num, label) in enumerate(items):
        col = i % 2
        row = i // 2
        left = left1 if col == 0 else left2
        top = top_start + row * row_h

        # number badge
        add_rect(slide, left, top + Inches(0.06), Inches(0.45), Inches(0.42), ACCENT)
        add_text_box(slide, num, left, top + Inches(0.03), Inches(0.45), Inches(0.48),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # label
        add_text_box(slide, label,
                     left + Inches(0.52), top + Inches(0.05), col_w - Inches(0.6), Inches(0.45),
                     font_size=14, color=PRIMARY, bold=False, align=PP_ALIGN.LEFT)


def slide_problem_statement(prs):
    """Slide 3 – Problem Statement"""
    slide = add_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    add_header_bar(slide, "Problem Statement", "Why do we need a Corporate Carpooling System?")
    add_footer(slide)

    problems = [
        "🚦  Urban Traffic Congestion — Thousands of employees driving separately causes severe road congestion around corporate campuses, increasing commute time by up to 45%.",
        "🌍  Environmental Impact — Individual cars emit 4–5× more CO₂ per commute than shared rides. Corporate fleets remain the single largest source of avoidable employee carbon footprint.",
        "💸  High Commute Costs — Fuel, parking, and vehicle maintenance cost employees ₹8,000–₹20,000 per month in metro cities, reducing effective take-home pay.",
        "🔒  Safety Concerns for Solo Commuters — Women employees and night-shift workers face significant safety risks during early-morning or late-night commutes.",
        "📭  No Centralized Platform — Organisations lack a dedicated system to coordinate carpooling; existing public apps (Ola, Uber) are not designed for closed corporate ecosystems.",
        "📊  No Visibility for HR/Admin — Companies cannot measure commute patterns, carbon footprint, or employee transport costs without a dedicated system.",
    ]

    top = Inches(1.4)
    for prob in problems:
        add_text_box(slide, prob,
                     Inches(0.45), top, Inches(12.4), Inches(0.82),
                     font_size=13.5, color=DARK_GRAY, wrap=True)
        top += Inches(0.9)


def slide_solution(prs):
    """Slide 4 – Proposed Solution"""
    slide = add_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    add_header_bar(slide, "Proposed Solution", "CarpoolHub — A Corporate-Grade Ride Sharing Platform")
    add_footer(slide)

    add_text_box(slide,
                 "A web-based corporate carpooling platform enabling same-organisation employees to offer and request shared rides, reducing individual car usage, CO₂ emissions, and commute costs.",
                 Inches(0.45), Inches(1.38), Inches(12.4), Inches(0.75),
                 font_size=14.5, color=PRIMARY, bold=True, wrap=True)

    objectives = [
        ("Core Objectives", [
            "✔  Enable drivers to post rides with route, timing, and seat count",
            "✔  Smart matching of passengers to nearby drivers using geo-spatial algorithms",
            "✔  Real-time ride tracking via GPS pings",
            "✔  In-app chat between driver and passengers",
            "✔  Integrated payment system (Razorpay)",
            "✔  Mutual rating system after ride completion",
        ]),
        ("Safety & Reliability", [
            "✔  SOS emergency alert with guardian contact notification",
            "✔  Backup driver mechanism — auto-activate if primary cancels",
            "✔  Gender-preference filtering for ride matching",
            "✔  Soft-delete and audit trails for data safety",
            "✔  Organisation-scoped access (employees see only org rides)",
            "✔  AI Agent for architecture & system Q&A (LangGraph)",
        ]),
    ]

    col_w = Inches(6.1)
    for ci, (heading, bullets) in enumerate(objectives):
        left = Inches(0.35) + ci * (col_w + Inches(0.35))
        top = Inches(2.22)
        add_rect(slide, left, top, col_w, Inches(0.4), ACCENT)
        add_text_box(slide, heading,
                     left + Inches(0.1), top + Inches(0.04), col_w - Inches(0.2), Inches(0.34),
                     font_size=13, bold=True, color=WHITE)
        add_rect(slide, left, top + Inches(0.4), col_w, Inches(4.65), WHITE,
                 line_color=ACCENT, line_width=Pt(1))
        body_top = top + Inches(0.5)
        for b in bullets:
            add_text_box(slide, b,
                         left + Inches(0.15), body_top, col_w - Inches(0.25), Inches(0.62),
                         font_size=12.5, color=DARK_GRAY)
            body_top += Inches(0.68)


def slide_architecture(prs):
    """Slide 5 – System Architecture"""
    slide = add_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    add_header_bar(slide, "System Architecture", "3-Tier Architecture with AI-Augmented Layer")
    add_footer(slide)

    # Architecture diagram
    svg_path = os.path.join(DIAGRAMS, "Architecture Diagram.svg")
    png_bytes = svg_to_png_bytes(svg_path)
    from PIL import Image
    img = Image.open(io.BytesIO(png_bytes))
    orig_w, orig_h = img.size
    aspect = orig_w / orig_h
    img_h = Inches(5.75)
    img_w = img_h * aspect
    if img_w > Inches(12.8):
        img_w = Inches(12.8)
        img_h = img_w / aspect
    left = (Inches(13.33) - img_w) / 2
    slide.shapes.add_picture(io.BytesIO(png_bytes), left, Inches(1.35), img_w, img_h)


def slide_tech_stack(prs):
    """Slide 6 – Technology Stack"""
    slide = add_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    add_header_bar(slide, "Technology Stack", "Tools, Frameworks & Services")
    add_footer(slide)

    layers = [
        ("Frontend", ACCENT, [
            "React 18 (Vite)",
            "React Router v6",
            "Context API (state)",
            "Axios (REST calls)",
            "CSS Modules + Tokens",
        ]),
        ("Backend", PRIMARY, [
            "Spring Boot 3 (Java 21)",
            "Spring Security + JWT",
            "Spring Data JPA",
            "Hibernate / Validation",
            "WebSocket (STOMP)",
        ]),
        ("Database", RGBColor(0x6C, 0x35, 0x8A), [
            "PostgreSQL 16",
            "PostGIS Extension",
            "ST_DWithin / ST_Distance",
            "GEOGRAPHY(Point, 4326)",
            "GEOMETRY(LineString, 4326)",
        ]),
        ("Integrations", GREEN, [
            "Razorpay (Payments)",
            "Claude AI (Anthropic)",
            "LangGraph (AI Agent)",
            "FAISS (Vector Store)",
            "JWT / BCrypt",
        ]),
        ("DevOps / Infra", ORANGE, [
            "Git + GitHub",
            "Maven (build)",
            "Vite (frontend build)",
            "Python 3.11 (AI Agent)",
            "PostgreSQL + PostGIS",
        ]),
    ]

    card_w = Inches(2.3)
    card_h = Inches(5.6)
    gap    = Inches(0.27)
    top    = Inches(1.35)
    left   = Inches(0.35)

    for li, (layer_name, layer_color, items) in enumerate(layers):
        lx = left + li * (card_w + gap)
        # header
        add_rect(slide, lx, top, card_w, Inches(0.42), layer_color)
        add_text_box(slide, layer_name,
                     lx + Inches(0.05), top + Inches(0.04), card_w - Inches(0.1), Inches(0.35),
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # body
        add_rect(slide, lx, top + Inches(0.42), card_w, card_h - Inches(0.42), WHITE,
                 line_color=layer_color, line_width=Pt(1.5))
        bt = top + Inches(0.55)
        for item in items:
            add_text_box(slide, f"•  {item}",
                         lx + Inches(0.12), bt, card_w - Inches(0.2), Inches(0.88),
                         font_size=12, color=DARK_GRAY, wrap=True)
            bt += Inches(0.95)


def slide_modules(prs):
    """Slide 7 – Core Modules"""
    slide = add_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    add_header_bar(slide, "Core Modules", "9 Functional Modules forming the complete system")
    add_footer(slide)

    modules = [
        ("👤 User Management",     "Registration, login, JWT auth, profile, soft-delete, org scoping"),
        ("🚗 Ride Management",      "Create/update/cancel rides, seat tracking, status lifecycle"),
        ("🔍 Matching Engine",      "Geo-spatial matching via ST_DWithin, route overlap, detour %, gender filter, rating sort"),
        ("📋 Request Handling",     "Passenger ride requests, approval/rejection by driver, request lifecycle"),
        ("📍 Live Tracking",        "GPS pings every N seconds, real-time location stream via WebSocket"),
        ("💬 Chat System",          "Ride-based messaging, read receipts, WebSocket delivery"),
        ("💳 Payment System",       "Razorpay integration, INITIATED → SUCCESS → FAILED → REFUNDED"),
        ("🛡 Safety (SOS)",         "One-tap SOS, guardian contact notification, incident storage"),
        ("🔄 Backup Driver",        "Priority backup drivers per ride, auto-activation on primary cancel"),
    ]

    col_w = Inches(6.1)
    col_gap = Inches(0.45)
    left1 = Inches(0.35)
    left2 = left1 + col_w + col_gap
    top_s = Inches(1.38)
    row_h = Inches(0.68)

    for i, (title, desc) in enumerate(modules):
        col = i % 2
        row = i // 2
        lx = left1 if col == 0 else left2
        ty = top_s + row * row_h
        add_rect(slide, lx, ty, col_w, row_h - Inches(0.06), WHITE,
                 line_color=ACCENT, line_width=Pt(0.8))
        add_text_box(slide, title,
                     lx + Inches(0.1), ty + Inches(0.05), col_w - Inches(0.2), Inches(0.3),
                     font_size=13, bold=True, color=PRIMARY)
        add_text_box(slide, desc,
                     lx + Inches(0.1), ty + Inches(0.32), col_w - Inches(0.2), Inches(0.3),
                     font_size=11, color=DARK_GRAY, wrap=True)

    # last module (9th) centred
    i = 8
    title, desc = modules[i]
    lx = left1 + (col_w + col_gap) / 2 - Inches(0.05)
    ty = top_s + 4 * row_h
    add_rect(slide, lx, ty, col_w, row_h - Inches(0.06), WHITE,
             line_color=ACCENT, line_width=Pt(0.8))
    add_text_box(slide, title,
                 lx + Inches(0.1), ty + Inches(0.05), col_w - Inches(0.2), Inches(0.3),
                 font_size=13, bold=True, color=PRIMARY)
    add_text_box(slide, desc,
                 lx + Inches(0.1), ty + Inches(0.32), col_w - Inches(0.2), Inches(0.3),
                 font_size=11, color=DARK_GRAY, wrap=True)


def slide_er_diagram(prs):
    return add_diagram_slide(
        prs,
        "Database Design — ER Diagram",
        "PostgreSQL 16 + PostGIS | 15 tables with geo-spatial support",
        "ER Diagram Final.svg",
        notes="Key tables: organisations, users, vehicles, routes, ride_schedules, ride_requests, ride_passengers, backup_rides, sos_incidents, transactions, ratings, chat_messages, ride_location_pings, ride_events, matching_logs"
    )


def slide_db_tables(prs):
    """Slide – Database Tables Summary"""
    slide = add_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    add_header_bar(slide, "Database Design — Key Tables",
                   "15 tables | PostGIS geo-spatial extensions | Audit triggers")
    add_footer(slide)

    tables = [
        ("organisations",        "Company registry; scopes all users & rides"),
        ("users",                "Auth, profile, is_online, is_deleted, org FK"),
        ("vehicles",             "Driver's vehicles; one active vehicle per driver"),
        ("routes",               "GEOMETRY(LineString) for driver's route path"),
        ("ride_schedules",       "Core ride entity; lifecycle CREATED→ACTIVE→STARTED→COMPLETED"),
        ("ride_requests",        "Passenger intent with GEOGRAPHY pickup/drop points"),
        ("ride_passengers",      "Confirmed riders per ride (ACTIVE/CANCELLED/COMPLETED)"),
        ("ride_location_pings",  "GPS pings for live tracking (GEOGRAPHY Point)"),
        ("backup_rides",         "Backup driver entries; PENDING→ACTIVATED→USED/EXPIRED"),
        ("chat_messages",        "Ride-scoped messages with is_read & sender_id"),
        ("sos_incidents",        "SOS triggers with location & status"),
        ("transactions",         "Razorpay payments per ride + user; full state machine"),
        ("ratings",              "Post-ride mutual ratings (driver↔passenger)"),
        ("ride_events",          "Full audit trail of every ride state change"),
        ("matching_logs",        "Detour %, proximity, match decision per request"),
    ]

    col_w1 = Inches(3.1)
    col_w2 = Inches(9.4)
    row_h  = Inches(0.395)
    lx     = Inches(0.35)
    ty     = Inches(1.38)

    # header row
    add_rect(slide, lx, ty, col_w1, row_h, PRIMARY)
    add_text_box(slide, "Table", lx + Inches(0.1), ty + Inches(0.04),
                 col_w1, row_h, font_size=12, bold=True, color=WHITE)
    add_rect(slide, lx + col_w1, ty, col_w2, row_h, PRIMARY)
    add_text_box(slide, "Purpose", lx + col_w1 + Inches(0.1), ty + Inches(0.04),
                 col_w2, row_h, font_size=12, bold=True, color=WHITE)
    ty += row_h

    for i, (tbl, desc) in enumerate(tables):
        bg = WHITE if i % 2 == 0 else RGBColor(0xE8, 0xF0, 0xFE)
        add_rect(slide, lx, ty, col_w1, row_h, bg,
                 line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
        add_text_box(slide, tbl, lx + Inches(0.08), ty + Inches(0.04),
                     col_w1 - Inches(0.1), row_h,
                     font_size=11, bold=True, color=PRIMARY)
        add_rect(slide, lx + col_w1, ty, col_w2, row_h, bg,
                 line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
        add_text_box(slide, desc, lx + col_w1 + Inches(0.08), ty + Inches(0.04),
                     col_w2 - Inches(0.1), row_h, font_size=11, color=DARK_GRAY)
        ty += row_h


def slide_use_case(prs):
    return add_diagram_slide(
        prs,
        "Use Case Diagram",
        "Actors: Driver, Passenger, Admin | 20+ use cases",
        "Use Case Diagram.svg",
    )


def slide_activity(prs):
    return add_diagram_slide(
        prs,
        "Activity Diagram — Overall Flow",
        "End-to-end ride lifecycle from login to ride completion",
        "Activity Diagram.svg",
    )


def slide_activity_booking(prs):
    return add_diagram_slide(
        prs,
        "Activity Diagram — Ride Booking Flow",
        "Passenger ride request → Matching → Acceptance → Tracking",
        "Ride Booking activity Diagram.svg",
    )


def slide_combined_activity(prs):
    return add_diagram_slide(
        prs,
        "Combined Activity Flow Diagram",
        "Unified view of Driver and Passenger parallel flows",
        "Combined activity flow diagram.svg",
    )


def slide_sequence(prs):
    return add_diagram_slide(
        prs,
        "Sequence Diagram",
        "System interactions between Frontend, Backend, DB and Payment Gateway",
        "Sequence Diagram Final.svg",
    )


def slide_backup_sequence(prs):
    return add_diagram_slide(
        prs,
        "Backup Driver — Sequence Diagram",
        "How the system activates a backup driver when primary cancels",
        "Backup Driver Sequnece diagram final.svg",
    )


def slide_backup_flow(prs):
    return add_diagram_slide(
        prs,
        "Backup Driver — Activation Flow",
        "Priority-based backup selection and state transitions",
        "Backup driver activation flow.svg",
    )


def slide_state_machine(prs):
    return add_diagram_slide(
        prs,
        "Ride Lifecycle — State Machine",
        "States: CREATED → ACTIVE → STARTED → COMPLETED | CANCELLED",
        "Ride lifecycle State machine.svg",
    )


def slide_matching_engine(prs):
    """Slide – Geo-Spatial Matching Engine"""
    slide = add_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    add_header_bar(slide, "Matching Engine — Core Algorithm",
                   "Geo-spatial ride matching using PostGIS")
    add_footer(slide)

    steps = [
        ("Step 1: Proximity Filter",
         "ST_DWithin(driver.start_location, passenger.pickup_location, radius)\n→ Find all active drivers within configurable radius (default 5 km)"),
        ("Step 2: Route Overlap Check",
         "ST_Distance(route.geometry, passenger.pickup) < threshold\n→ Ensure driver's route passes near passenger pickup & drop"),
        ("Step 3: Detour Calculation",
         "detour% = (new_distance − original_distance) / original_distance × 100\n→ Filter out rides exceeding driver's detourLimitPercent"),
        ("Step 4: Constraint Filters",
         "available_seats ≥ 1\ngender_preference matches passenger gender\nride status = ACTIVE"),
        ("Step 5: Ranking & Result",
         "Sort by: distance ASC → rating DESC → departure time ASC\n→ Return top N matches to passenger"),
    ]

    top = Inches(1.4)
    for i, (title, desc) in enumerate(steps):
        add_rect(slide, Inches(0.35), top, Inches(0.52), Inches(1.05),
                 ACCENT if i % 2 == 0 else PRIMARY)
        add_text_box(slide, str(i + 1),
                     Inches(0.35), top + Inches(0.25), Inches(0.52), Inches(0.52),
                     font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(0.94), top, Inches(11.94), Inches(1.05), WHITE,
                 line_color=ACCENT, line_width=Pt(0.8))
        add_text_box(slide, title,
                     Inches(1.08), top + Inches(0.05), Inches(11.7), Inches(0.35),
                     font_size=13, bold=True, color=PRIMARY)
        add_text_box(slide, desc,
                     Inches(1.08), top + Inches(0.4), Inches(11.7), Inches(0.58),
                     font_size=11.5, color=DARK_GRAY, wrap=True)
        top += Inches(1.14)


def slide_backend_progress(prs):
    """Slide – Backend Implementation Progress"""
    slide = add_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    add_header_bar(slide, "Backend Implementation — Progress",
                   "Spring Boot 3 | 16 REST Controllers | JWT Auth | Razorpay")
    add_footer(slide)

    controllers = [
        "AuthController",  "UserController",       "VehicleController",
        "OrganisationController", "RideScheduleController", "RideRequestController",
        "RidePassengerController", "MatchingController", "TrackingController",
        "ChatController",  "PaymentController",    "RatingController",
        "SosController",   "BackupRideController", "GuardianContactController",
        "RouteWaypointController",
    ]

    done_items = [
        "✅  JWT authentication & Spring Security filter chain",
        "✅  User management with soft-delete & org scoping",
        "✅  Ride lifecycle: CREATED → ACTIVE → STARTED → COMPLETED",
        "✅  Geo-spatial matching engine with PostGIS",
        "✅  Backup driver system with priority-based activation",
        "✅  Razorpay payment integration (initiate / confirm / refund)",
        "✅  SOS system with guardian contact notifications",
        "✅  Live GPS tracking (location ping storage)",
        "✅  Chat system (REST; WebSocket planned)",
        "✅  Rating system (driver ↔ passenger)",
        "✅  Vehicle & Organisation management APIs",
        "✅  Role-based access guard (owner-only delete)",
        "✅  Passenger status check before joining rides",
        "✅  Gender preference filtering in matching",
    ]

    # Controllers grid
    add_text_box(slide, "REST Controllers (16):",
                 Inches(0.35), Inches(1.38), Inches(5.5), Inches(0.35),
                 font_size=13, bold=True, color=PRIMARY)

    col_w = Inches(1.65)
    gap   = Inches(0.08)
    lx    = Inches(0.35)
    ty    = Inches(1.75)
    for i, ctrl in enumerate(controllers):
        row = i // 3
        col = i % 3
        cx = lx + col * (col_w + gap)
        cy = ty + row * Inches(0.42)
        add_rect(slide, cx, cy, col_w, Inches(0.38), ACCENT)
        add_text_box(slide, ctrl.replace("Controller", ""),
                     cx + Inches(0.05), cy + Inches(0.04), col_w - Inches(0.1), Inches(0.32),
                     font_size=10.5, color=WHITE, align=PP_ALIGN.CENTER)

    # Done items on right
    add_text_box(slide, "Implemented Features:",
                 Inches(5.6), Inches(1.38), Inches(7.3), Inches(0.35),
                 font_size=13, bold=True, color=PRIMARY)
    ty2 = Inches(1.75)
    for item in done_items:
        add_text_box(slide, item,
                     Inches(5.6), ty2, Inches(7.3), Inches(0.41),
                     font_size=11, color=DARK_GRAY)
        ty2 += Inches(0.43)


def slide_frontend_progress(prs):
    """Slide – Frontend Progress"""
    slide = add_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    add_header_bar(slide, "Frontend Implementation — Progress",
                   "React 18 + Vite | Context API | Axios | Responsive UI")
    add_footer(slide)

    done = [
        "✅  Project scaffolded with Vite + React 18",
        "✅  Authentication screens (Login, Register)",
        "✅  Context API for auth state management",
        "✅  Axios API client with JWT interceptor",
        "✅  Design system tokens (colours, spacing, typography)",
        "✅  Route setup with React Router v6",
        "✅  Component library structure in place",
        "✅  Admin module scaffolded",
        "✅  Screens directory organised by role",
    ]

    planned = [
        "🔲  Dashboard screen (Driver & Passenger views)",
        "🔲  Ride scheduling form with map picker",
        "🔲  Ride request & matching results screen",
        "🔲  Live tracking map (Google Maps / Leaflet)",
        "🔲  Chat interface (WebSocket integration)",
        "🔲  Payment flow screens (Razorpay SDK)",
        "🔲  SOS trigger button (persistent on ride screen)",
        "🔲  Rating & review screen post-ride",
        "🔲  Admin dashboard (ride analytics, user mgmt)",
    ]

    col_w = Inches(6.1)
    gap   = Inches(0.45)
    lx1   = Inches(0.35)
    lx2   = lx1 + col_w + gap

    for lx, heading, items, hcolor in [
        (lx1, "Completed", done, GREEN),
        (lx2, "In Progress / Planned", planned, ORANGE),
    ]:
        add_rect(slide, lx, Inches(1.38), col_w, Inches(0.42), hcolor)
        add_text_box(slide, heading,
                     lx + Inches(0.1), Inches(1.4), col_w - Inches(0.2), Inches(0.38),
                     font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        ty = Inches(1.85)
        for item in items:
            add_text_box(slide, item,
                         lx + Inches(0.1), ty, col_w - Inches(0.2), Inches(0.52),
                         font_size=12, color=DARK_GRAY)
            ty += Inches(0.56)


def slide_ai_agent(prs):
    """Slide – AI Agent Module"""
    slide = add_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    add_header_bar(slide, "AI Agent Module — CarpoolHub Query System",
                   "LangGraph + Claude AI + FAISS | Token-efficient architecture Q&A")
    add_footer(slide)

    # Architecture description
    add_text_box(slide,
                 "An intelligent LangGraph-based AI agent that routes architecture queries to relevant knowledge sources, compresses context for token efficiency, and answers using Claude (Anthropic).",
                 Inches(0.45), Inches(1.4), Inches(12.4), Inches(0.65),
                 font_size=13.5, color=PRIMARY, bold=False, wrap=True)

    components = [
        ("Intent Classifier",   "classifier.py", "Routes queries to relevant sub-agents (DB, flows, services, architecture)"),
        ("Vector Retriever",    "retriever.py",  "FAISS vector store over SVG diagram text & JSON data; semantic search"),
        ("Context Compressor",  "compressor.py", "Minimises tokens sent to Claude by pruning irrelevant context"),
        ("LangGraph Nodes",     "nodes.py",      "Structured + vector retrieval nodes wired into directed graph"),
        ("Orchestrator",        "graph.py",       "LangGraph StateGraph managing query → retrieve → compress → answer"),
        ("Knowledge Base",      "data/",          "Structured JSON (DB schema, flows, services) + raw SVG text of all diagrams"),
    ]

    top = Inches(2.12)
    card_h = Inches(0.8)
    for i, (name, file_, desc) in enumerate(components):
        col = i % 2
        row = i // 2
        lx = Inches(0.35) + col * Inches(6.6)
        ty = top + row * (card_h + Inches(0.1))
        add_rect(slide, lx, ty, Inches(6.4), card_h, WHITE,
                 line_color=ACCENT, line_width=Pt(1))
        add_rect(slide, lx, ty, Inches(1.8), card_h, PRIMARY)
        add_text_box(slide, name,
                     lx + Inches(0.06), ty + Inches(0.1), Inches(1.68), Inches(0.3),
                     font_size=11, bold=True, color=WHITE, wrap=True)
        add_text_box(slide, file_,
                     lx + Inches(0.06), ty + Inches(0.42), Inches(1.68), Inches(0.28),
                     font_size=9.5, color=ACCENT, italic=True)
        add_text_box(slide, desc,
                     lx + Inches(1.88), ty + Inches(0.14), Inches(4.4), Inches(0.55),
                     font_size=11.5, color=DARK_GRAY, wrap=True)

    add_text_box(slide,
                 "Agent Flow:  User Query → Intent Classification → Structured + Vector Retrieval → Context Compression → Claude AI → Answer",
                 Inches(0.45), Inches(6.72), Inches(12.4), Inches(0.42),
                 font_size=12.5, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)


def slide_progress_summary(prs):
    """Slide – Progress Summary"""
    slide = add_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    add_header_bar(slide, "Progress Summary", "What has been completed so far")
    add_footer(slide)

    categories = [
        ("Database & Schema", GREEN, [
            "PostgreSQL 16 + PostGIS set up",
            "15 production-grade tables designed",
            "Migration scripts (V1, V2, V3) applied",
            "Geo-spatial indexes (GIST) created",
            "Audit triggers on all tables",
            "Constraints: seats ≥ 0, unique active ride/driver",
        ]),
        ("Backend (Spring Boot)", GREEN, [
            "16 REST Controllers implemented",
            "JWT authentication & security filter",
            "Ride lifecycle state machine",
            "Geo-spatial matching engine",
            "Razorpay payment integration",
            "Backup driver & SOS systems",
        ]),
        ("Frontend (React)", ORANGE, [
            "Project scaffolded & routing set up",
            "Auth screens (Login/Register) done",
            "Context API & Axios client ready",
            "Design tokens defined",
            "Admin & screens structure in place",
            "Remaining screens under development",
        ]),
        ("AI Agent & Docs", GREEN, [
            "LangGraph AI agent implemented",
            "FAISS vector store created",
            "All diagrams text-encoded",
            "Architecture diagrams (10 SVGs)",
            "ER, Sequence, Activity, State Machine",
            "OpenAPI / Swagger — in progress",
        ]),
    ]

    col_w = Inches(3.1)
    gap   = Inches(0.15)
    top   = Inches(1.38)
    for ci, (title, color, items) in enumerate(categories):
        lx = Inches(0.35) + ci * (col_w + gap)
        add_rect(slide, lx, top, col_w, Inches(0.42), color)
        add_text_box(slide, title,
                     lx + Inches(0.08), top + Inches(0.05), col_w - Inches(0.15), Inches(0.34),
                     font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, lx, top + Inches(0.42), col_w, Inches(5.55), WHITE,
                 line_color=color, line_width=Pt(1.2))
        ty = top + Inches(0.54)
        for item in items:
            add_text_box(slide, f"• {item}",
                         lx + Inches(0.1), ty, col_w - Inches(0.18), Inches(0.8),
                         font_size=11, color=DARK_GRAY, wrap=True)
            ty += Inches(0.85)


def slide_future_roadmap(prs):
    """Slide – Future Roadmap"""
    slide = add_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    add_header_bar(slide, "Future Roadmap", "What we plan to complete before final submission")
    add_footer(slide)

    phases = [
        ("Phase 1 — Immediate (Next 2 Weeks)", ACCENT, [
            "Complete all React screens: Dashboard, Ride Scheduling, Matching Results",
            "Integrate live map tracking (Leaflet.js / Google Maps)",
            "Connect frontend to all backend APIs",
            "Implement WebSocket-based real-time chat",
            "Add Razorpay payment flow in frontend",
        ]),
        ("Phase 2 — Short Term (Weeks 3–5)", PRIMARY, [
            "SOS button with real-time guardian notification",
            "Swagger / OpenAPI documentation for all 40+ endpoints",
            "Backup driver activation UI",
            "Rating & review screen post-ride",
            "Admin dashboard (ride analytics, user management)",
        ]),
        ("Phase 3 — Final Polish (Week 6)", GREEN, [
            "End-to-end integration testing (Playwright / Jest)",
            "Performance optimisation (DB query plans, indexing review)",
            "Security audit (OWASP top 10, SQL injection checks)",
            "Deployment setup (Docker / cloud hosting)",
            "Final dissertation report & demo preparation",
        ]),
        ("Future Scope (Post-MTech)", ORANGE, [
            "Native mobile apps (Android & iOS — React Native)",
            "ML-based route prediction & proactive ride suggestions",
            "Corporate analytics dashboard (CO₂ saved, cost reports)",
            "Multi-organisation / SaaS model",
            "Driver earnings & incentive system",
        ]),
    ]

    top = Inches(1.38)
    for pi, (phase_name, color, items) in enumerate(phases):
        ty = top + pi * Inches(1.55)
        add_rect(slide, Inches(0.35), ty, Inches(3.5), Inches(1.45), color)
        add_text_box(slide, phase_name,
                     Inches(0.42), ty + Inches(0.2), Inches(3.35), Inches(1.05),
                     font_size=12, bold=True, color=WHITE, wrap=True, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(3.88), ty, Inches(9.1), Inches(1.45), WHITE,
                 line_color=color, line_width=Pt(1.2))
        tx = Inches(3.98)
        item_w = Inches(4.3)
        item_h = Inches(0.26)
        for ii, item in enumerate(items):
            ix = ii % 2
            iy = ii // 2
            add_text_box(slide, f"→  {item}",
                         tx + ix * item_w, ty + Inches(0.08) + iy * Inches(0.46),
                         item_w, item_h + Inches(0.16),
                         font_size=11, color=DARK_GRAY, wrap=True)


def slide_thankyou(prs):
    """Slide – Thank You"""
    slide = add_slide(prs)
    fill_bg(slide, PRIMARY)

    add_rect(slide, 0, Inches(3.3), Inches(13.33), Inches(0.06), ACCENT)

    add_text_box(slide, "Thank You",
                 Inches(1.0), Inches(1.5), Inches(11.33), Inches(1.0),
                 font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Questions & Discussion Welcome",
                 Inches(1.0), Inches(2.55), Inches(11.33), Inches(0.55),
                 font_size=22, color=ACCENT, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(3.5), Inches(3.55), Inches(6.33), Inches(2.1),
             RGBColor(0x0D, 0x22, 0x4A))
    add_text_box(slide, "Prashant Sopan Berad",
                 Inches(3.65), Inches(3.65), Inches(6.0), Inches(0.45),
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "MTech (Software Engineering)\nBITS Pilani",
                 Inches(3.65), Inches(4.1), Inches(6.0), Inches(0.6),
                 font_size=14, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Corporate Car Pooling System | May 2025",
                 Inches(3.65), Inches(4.72), Inches(6.0), Inches(0.42),
                 font_size=12, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "GitHub: github.com/infotechmaniac17/Corporate-Car-Pooling-System",
                 Inches(1.0), Inches(5.9), Inches(11.33), Inches(0.38),
                 font_size=13, color=RGBColor(0x88, 0xAA, 0xCC),
                 align=PP_ALIGN.CENTER, italic=True)

    add_footer(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)   # widescreen 16:9
    prs.slide_height = Inches(7.5)

    print("Building slides...")
    slide_title(prs);              print("  1/20  Title")
    slide_agenda(prs);             print("  2/20  Agenda")
    slide_problem_statement(prs);  print("  3/20  Problem Statement")
    slide_solution(prs);           print("  4/20  Proposed Solution")
    slide_architecture(prs);       print("  5/20  Architecture Diagram")
    slide_tech_stack(prs);         print("  6/20  Tech Stack")
    slide_modules(prs);            print("  7/20  Core Modules")
    slide_er_diagram(prs);         print("  8/20  ER Diagram")
    slide_db_tables(prs);          print("  9/20  DB Tables")
    slide_use_case(prs);           print(" 10/20  Use Case Diagram")
    slide_activity(prs);           print(" 11/20  Activity Diagram")
    slide_activity_booking(prs);   print(" 12/20  Ride Booking Activity")
    slide_combined_activity(prs);  print(" 13/20  Combined Activity Flow")
    slide_sequence(prs);           print(" 14/20  Sequence Diagram")
    slide_backup_sequence(prs);    print(" 15/20  Backup Driver Sequence")
    slide_backup_flow(prs);        print(" 16/20  Backup Driver Flow")
    slide_state_machine(prs);      print(" 17/20  State Machine")
    slide_matching_engine(prs);    print(" 18/20  Matching Engine")
    slide_backend_progress(prs);   print(" 19/20  Backend Progress")
    slide_frontend_progress(prs);  print(" 20/20  Frontend Progress")
    slide_ai_agent(prs);           print(" 21/22  AI Agent")
    slide_progress_summary(prs);   print(" 22/23  Progress Summary")
    slide_future_roadmap(prs);     print(" 23/24  Future Roadmap")
    slide_thankyou(prs);           print(" 24/24  Thank You")

    out_path = "/home/user/Corporate-Car-Pooling-System/Corporate_Carpooling_Dissertation.pptx"
    prs.save(out_path)
    print(f"\nSaved: {out_path}")


if __name__ == "__main__":
    main()
